#!/usr/bin/env python3
"""
Generate PowerPoint from Executive Leadership Lab Module 1 content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors from design system
CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)
TEAL = RGBColor(0x2A, 0x7B, 0x8C)
WARM_GOLD = RGBColor(0xD4, 0xA5, 0x74)
CREAM = RGBColor(0xFA, 0xF8, 0xF5)
TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66)

def add_title_slide(prs, title, subtitle, section_label=None):
    """Add a dark title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = CHARCOAL
    background.line.fill.background()

    # Section label
    if section_label:
        label_box = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(8.5), Inches(0.4))
        label_tf = label_box.text_frame
        label_p = label_tf.paragraphs[0]
        label_p.text = section_label.upper()
        label_p.font.size = Pt(11)
        label_p.font.color.rgb = TEAL
        label_p.font.bold = True
        label_p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(8.5), Inches(1.5))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.color.rgb = CREAM
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(18)
        sub_p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        sub_p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, section_label, title, body_text=None, bullets=None):
    """Add a content slide with cream background"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Cream background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = CREAM
    background.line.fill.background()

    # Section label
    label_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(0.4))
    label_tf = label_box.text_frame
    label_p = label_tf.paragraphs[0]
    label_p.text = section_label.upper()
    label_p.font.size = Pt(10)
    label_p.font.color.rgb = TEAL
    label_p.font.bold = True
    label_p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1), Inches(8.5), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)
    title_p.font.color.rgb = CHARCOAL
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER

    # Body text
    if body_text:
        body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_p = body_tf.paragraphs[0]
        body_p.text = body_text
        body_p.font.size = Pt(16)
        body_p.font.color.rgb = TEXT_LIGHT
        body_p.alignment = PP_ALIGN.CENTER

    # Bullets
    if bullets:
        y_pos = 2.8 if body_text else 2
        bullet_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(7), Inches(4))
        bullet_tf = bullet_box.text_frame
        bullet_tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = bullet_tf.paragraphs[0]
            else:
                p = bullet_tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_LIGHT
            p.space_after = Pt(12)
            p.level = 0

    return slide

def add_quote_slide(prs, quote_text):
    """Add a quote slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Warm background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xF5, 0xF2, 0xED)
    background.line.fill.background()

    # Quote mark
    quote_mark = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(1), Inches(1))
    qm_tf = quote_mark.text_frame
    qm_p = qm_tf.paragraphs[0]
    qm_p.text = '"'
    qm_p.font.size = Pt(72)
    qm_p.font.color.rgb = WARM_GOLD

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
    quote_tf = quote_box.text_frame
    quote_tf.word_wrap = True
    quote_p = quote_tf.paragraphs[0]
    quote_p.text = quote_text
    quote_p.font.size = Pt(24)
    quote_p.font.color.rgb = CHARCOAL
    quote_p.font.italic = True
    quote_p.alignment = PP_ALIGN.CENTER

    return slide

def add_cards_slide(prs, section_label, title, subtitle, cards):
    """Add a slide with card boxes"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Cream background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = CREAM
    background.line.fill.background()

    # Section label
    label_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(8.5), Inches(0.3))
    label_tf = label_box.text_frame
    label_p = label_tf.paragraphs[0]
    label_p.text = section_label.upper()
    label_p.font.size = Pt(10)
    label_p.font.color.rgb = TEAL
    label_p.font.bold = True
    label_p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(8.5), Inches(0.6))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(28)
    title_p.font.color.rgb = CHARCOAL
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(8), Inches(0.5))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(14)
        sub_p.font.color.rgb = TEXT_LIGHT
        sub_p.alignment = PP_ALIGN.CENTER

    # Cards - arrange in grid
    num_cards = len(cards)
    cols = 3 if num_cards > 4 else 2
    card_width = 2.8 if cols == 3 else 4
    start_x = 0.5 if cols == 3 else 1

    for i, (card_title, card_text) in enumerate(cards):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_width + 0.2)
        y = 2.1 + row * 1.5

        # Card box
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.color.rgb = RGBColor(0xE5, 0xE0, 0xD8)

        # Card title
        ct_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(card_width - 0.3), Inches(0.4))
        ct_tf = ct_box.text_frame
        ct_p = ct_tf.paragraphs[0]
        ct_p.text = card_title
        ct_p.font.size = Pt(12)
        ct_p.font.color.rgb = CHARCOAL
        ct_p.font.bold = True

        # Card text
        cb_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.5), Inches(card_width - 0.3), Inches(0.8))
        cb_tf = cb_box.text_frame
        cb_tf.word_wrap = True
        cb_p = cb_tf.paragraphs[0]
        cb_p.text = card_text
        cb_p.font.size = Pt(10)
        cb_p.font.color.rgb = TEXT_LIGHT

    return slide

def add_numbered_list_slide(prs, section_label, title, subtitle, items):
    """Add a slide with numbered items"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Cream background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = CREAM
    background.line.fill.background()

    # Section label
    label_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(8.5), Inches(0.3))
    label_tf = label_box.text_frame
    label_p = label_tf.paragraphs[0]
    label_p.text = section_label.upper()
    label_p.font.size = Pt(10)
    label_p.font.color.rgb = TEAL
    label_p.font.bold = True
    label_p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(8.5), Inches(0.6))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(28)
    title_p.font.color.rgb = CHARCOAL
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(1.35), Inches(8), Inches(0.5))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(14)
        sub_p.font.color.rgb = TEXT_LIGHT
        sub_p.alignment = PP_ALIGN.CENTER

    # Numbered items
    y_start = 2.0
    for i, (item_title, item_text) in enumerate(items):
        y = y_start + i * 0.9

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(y), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = WARM_GOLD
        circle.line.fill.background()

        # Number text
        num_box = slide.shapes.add_textbox(Inches(1.2), Inches(y + 0.05), Inches(0.35), Inches(0.3))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = str(i + 1)
        num_p.font.size = Pt(12)
        num_p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        num_p.font.bold = True
        num_p.alignment = PP_ALIGN.CENTER

        # Item title
        it_box = slide.shapes.add_textbox(Inches(1.7), Inches(y), Inches(6.5), Inches(0.35))
        it_tf = it_box.text_frame
        it_p = it_tf.paragraphs[0]
        it_p.text = item_title
        it_p.font.size = Pt(14)
        it_p.font.color.rgb = CHARCOAL
        it_p.font.bold = True

        # Item text
        ib_box = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.35), Inches(6.5), Inches(0.5))
        ib_tf = ib_box.text_frame
        ib_tf.word_wrap = True
        ib_p = ib_tf.paragraphs[0]
        ib_p.text = item_text
        ib_p.font.size = Pt(11)
        ib_p.font.color.rgb = TEXT_LIGHT

    return slide

def add_two_column_slide(prs, section_label, title, subtitle, left_title, left_items, right_title, right_items):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Cream background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = CREAM
    background.line.fill.background()

    # Section label
    label_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(8.5), Inches(0.3))
    label_tf = label_box.text_frame
    label_p = label_tf.paragraphs[0]
    label_p.text = section_label.upper()
    label_p.font.size = Pt(10)
    label_p.font.color.rgb = TEAL
    label_p.font.bold = True
    label_p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(8.5), Inches(0.6))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(28)
    title_p.font.color.rgb = CHARCOAL
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(1.35), Inches(8), Inches(0.5))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(14)
        sub_p.font.color.rgb = TEXT_LIGHT
        sub_p.alignment = PP_ALIGN.CENTER

    # Left column header
    lh_box = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(4), Inches(0.4))
    lh_tf = lh_box.text_frame
    lh_p = lh_tf.paragraphs[0]
    lh_p.text = left_title
    lh_p.font.size = Pt(14)
    lh_p.font.color.rgb = CHARCOAL
    lh_p.font.bold = True

    # Left column line
    left_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.4), Inches(3.8), Inches(0.03))
    left_line.fill.solid()
    left_line.fill.fore_color.rgb = CHARCOAL
    left_line.line.fill.background()

    # Left column items
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.55), Inches(4), Inches(2.5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_tf.paragraphs[0]
        else:
            p = left_tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    # Right column header
    rh_box = slide.shapes.add_textbox(Inches(5.25), Inches(2), Inches(4), Inches(0.4))
    rh_tf = rh_box.text_frame
    rh_p = rh_tf.paragraphs[0]
    rh_p.text = right_title
    rh_p.font.size = Pt(14)
    rh_p.font.color.rgb = CHARCOAL
    rh_p.font.bold = True

    # Right column line
    right_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.25), Inches(2.4), Inches(3.8), Inches(0.03))
    right_line.fill.solid()
    right_line.fill.fore_color.rgb = CHARCOAL
    right_line.line.fill.background()

    # Right column items
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(2.55), Inches(4), Inches(2.5))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_tf.paragraphs[0]
        else:
            p = right_tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    return slide


def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # Slide 1: Title
    add_title_slide(prs,
        "Break Through\nCareer Plateaus",
        "Transform stagnation into strategic momentum. Build the clarity, confidence, and visibility to reach your next level.",
        "Module 1 — January"
    )

    # Slide 2: Welcome
    slide = add_content_slide(prs,
        "Welcome",
        "You're Not Stuck.\nYou're Ready to Elevate.",
        "If you've been doing the work but not getting the recognition, moving but not advancing, or leading but feeling invisible—this module is your breakthrough moment."
    )

    # Slide 3: Objectives
    add_numbered_list_slide(prs,
        "Objectives",
        "What You'll Gain",
        None,
        [
            ("Diagnose Your Plateau", "Identify exactly what type of plateau you're experiencing and the root causes holding you back"),
            ("Build Your Growth Roadmap", "Create a strategic 90-day plan with clear milestones and accountability checkpoints"),
            ("Take Your First Moves", "Leave with three concrete actions you can execute in the next 14 days to build momentum")
        ]
    )

    # Slide 4: The Challenge
    add_cards_slide(prs,
        "The Challenge",
        "Working Hard ≠ Moving Up",
        "You've been told to put your head down and deliver. But excellence alone doesn't create advancement.",
        [
            ("The Visibility Gap", "Your work speaks for itself—but it's speaking to the wrong people, or not at all."),
            ("The Hustle Trap", "More work won't get you promoted. Strategic positioning will."),
            ("The Access Problem", "Without sponsors in the room, opportunity finds someone else.")
        ]
    )

    # Slide 5: Framework
    add_cards_slide(prs,
        "Framework",
        "The Three Career Plateaus",
        "Every plateau falls into one of three categories. Understanding yours is the first step to breaking through.",
        [
            ("Skill Plateau", "You've maxed out your current capabilities. Growth requires new skills, experiences, or credentials."),
            ("Role Plateau", "Your position has a ceiling. No room to grow, limited scope, or organizational constraints."),
            ("Identity Plateau", "Internal barriers—imposter syndrome, perfectionism, or playing small—are holding you back.")
        ]
    )

    # Slide 6: Skill Plateau
    add_two_column_slide(prs,
        "Deep Dive",
        "Skill Plateau",
        "When your expertise brought you here, but won't take you there.",
        "Signs You're Here",
        ["Known as \"the expert\" but passed over for leadership", "Your strengths feel outdated or too narrow", "Lack of exposure to strategic work", "Comfort zone has become a ceiling"],
        "The Breakthrough",
        ["Identify 2-3 skills needed for your next role", "Seek stretch assignments, not just promotions", "Build cross-functional exposure", "Reposition your brand around future value"]
    )

    # Slide 7: Role Plateau
    add_two_column_slide(prs,
        "Deep Dive",
        "Role Plateau",
        "When the role itself has no runway—not your performance.",
        "Signs You're Here",
        ["Your manager has nowhere to go either", "The org is flat or restructuring", "Budget constraints limit new positions", "You've outgrown the scope but can't expand"],
        "The Breakthrough",
        ["Explore lateral moves with growth potential", "Consider adjacent departments or functions", "Build visibility outside your team", "Know when it's time to pivot externally"]
    )

    # Slide 8: Identity Plateau
    add_two_column_slide(prs,
        "Deep Dive",
        "Identity Plateau",
        "When the barrier isn't skill or opportunity—it's self-perception.",
        "Signs You're Here",
        ["Imposter syndrome despite track record", "Perfectionism delays action", "Fear of visibility or failure", "Shrinking to fit, not standing out"],
        "The Breakthrough",
        ["Challenge the internal narrative", "Take \"messy action\"—done beats perfect", "Build evidence of your capability", "Step into visibility before you feel ready"]
    )

    # Slide 9: Quote
    add_quote_slide(prs,
        "The ceiling you're hitting isn't made of glass. It's made of stories you've been told—and stories you tell yourself."
    )

    # Slide 10: Root Causes
    add_cards_slide(prs,
        "Root Causes",
        "What's Driving Your Plateau?",
        None,
        [
            ("Unclear Target", "No defined \"next best role\" or direction makes it impossible to move strategically"),
            ("Unclear Value Proposition", "Can't articulate your unique value or what you bring that others don't"),
            ("Low Visibility", "Working hard but decision-makers don't know your name or your impact"),
            ("Weak Sponsorship", "No one with power is advocating for you when opportunities arise"),
            ("Over-Functioning", "So busy doing everyone's work that there's no bandwidth for your growth"),
            ("Fear and Avoidance", "Avoiding the conversations, risks, or visibility that advancement requires")
        ]
    )

    # Slide 11: Self-Assessment
    add_content_slide(prs,
        "Self-Assessment",
        "Plateau Diagnostic",
        "Time to get honest about where you are. Use the Plateau Diagnostic worksheet to identify your primary plateau type and top drivers.",
        ["Part 1: Check all statements that feel true right now. Be honest—this is for your eyes only.", "Part 2: Notice which category has the most checkmarks—that's your primary plateau zone.", "Time: 5–7 minutes"]
    )

    # Slide 12: The Stuck Loop
    add_content_slide(prs,
        "Pattern Recognition",
        "The Stuck Loop",
        "Most plateaus follow a predictable pattern. Breaking free starts with recognizing yours.",
        ["Trigger → Thought → Emotion → Behavior → Outcome → (Repeat)", "", "Example: Asked to present → \"I'll mess this up\" → Anxiety → Over-prepare, speak fast → \"Fine\" but invisible → Reinforces doubt"]
    )

    # Slide 13: Map Your Stuck Loop
    add_numbered_list_slide(prs,
        "Reflection",
        "Map Your Stuck Loop",
        "Using the Stuck Loop worksheet, identify the pattern that keeps repeating. Then design a new choice.",
        [
            ("Identify Your Loop", "What trigger, thought, emotion, behavior, and outcome keep repeating?"),
            ("Design Your New Choice", "What new thought or behavior could interrupt this pattern?"),
            ("Name Your Support", "What reminder, tool, or accountability would help you make a different choice?")
        ]
    )

    # Slide 14: Career Growth Roadmap
    add_numbered_list_slide(prs,
        "Strategic Planning",
        "Your Career Growth Roadmap",
        "A structured approach to move from where you are to where you want to be. Follow these six steps in order.",
        [
            ("North Star", "Define your \"next best role\" and what success looks like"),
            ("Current Reality", "Honest assessment of where you are and what's blocking progress"),
            ("Strengths & Gaps", "What to leverage and what to develop for the next level"),
            ("90-Day Goals", "2-3 focused objectives with clear success measures"),
            ("First Moves", "Concrete actions for the next 14 days to build momentum"),
            ("Accountability", "Who will hold you to this and how often will you check in")
        ]
    )

    # Slide 15: Build Your Roadmap
    add_two_column_slide(prs,
        "Deep Work",
        "Build Your Roadmap",
        "Work through each section of the Career Growth Roadmap. Be specific—this document guides your year.",
        "Key Questions",
        ["What is my next best role?", "What must be true in 6-12 months?", "What are my top 3 strengths?", "What 2 gaps must I close?", "What 3 actions start this week?"],
        "Guidance",
        ["Be specific—\"get promoted\" isn't a goal", "Make it measurable—how will you know?", "Focus on what you can influence", "Name real barriers—they're data, not excuses"]
    )

    # Slide 16: Map Your Network
    add_cards_slide(prs,
        "Relationship Capital",
        "Map Your Network",
        "Before building, take stock. Who currently influences your growth? This baseline reveals gaps and opportunities.",
        [
            ("Manager / Direct Leader", "Who directly influences your development and reviews?"),
            ("Key Stakeholders", "2-3 people whose opinion shapes your reputation"),
            ("Potential Sponsors", "Who has power and might advocate for you behind closed doors?"),
            ("Mentors", "Who provides guidance, wisdom, or career advice?"),
            ("Peer Allies", "Trusted peers you collaborate with and lean on"),
            ("External Community", "Networks and groups outside your organization")
        ]
    )

    # Slide 17: 14-Day Sprint
    add_numbered_list_slide(prs,
        "Momentum",
        "The 14-Day Sprint",
        "Momentum starts with small, visible moves. Your first actions should be achievable wins that create forward motion.",
        [
            ("Schedule Your First Moves", "Block time for the 3 actions from your roadmap—make them non-negotiable"),
            ("Track Evidence of Impact", "Document wins, feedback, and visibility moments—you'll need this later"),
            ("Schedule One Strategic Conversation", "A manager check-in, sponsor touchpoint, or stakeholder meeting")
        ]
    )

    # Slide 18: Quote
    add_quote_slide(prs,
        "You don't have to be ready. You have to be willing to start before you feel prepared."
    )

    # Slide 19: Key Takeaways
    add_numbered_list_slide(prs,
        "Summary",
        "Key Takeaways",
        None,
        [
            ("Plateaus Are Diagnosable", "Skill, Role, or Identity—knowing your type unlocks the right strategy"),
            ("Hard Work Alone Won't Get You There", "Strategic positioning, visibility, and sponsorship matter as much as performance"),
            ("Patterns Repeat Until Interrupted", "Your stuck loop will continue until you make a different choice"),
            ("Momentum Beats Planning", "14 days of action creates more clarity than 90 days of thinking")
        ]
    )

    # Slide 20: Homework
    add_numbered_list_slide(prs,
        "Before Next Session",
        "Your Homework",
        None,
        [
            ("Complete Your Career Growth Roadmap", "If you didn't finish in session, complete all sections this week"),
            ("Execute Your First 3 Actions", "Complete the first moves from your 14-day momentum planner"),
            ("Schedule One Conversation", "Book a meeting with your manager, sponsor, or key stakeholder"),
            ("Document One Win", "Capture at least one accomplishment or impact moment this period")
        ]
    )

    # Slide 21: Conclusion
    add_title_slide(prs,
        "Your Breakthrough\nStarts Now",
        "You have the diagnosis. You have the roadmap. You have the first moves.\nThe only thing left is action.\n\nNext Module: Strategic Positioning for Promotion — February",
        "Executive Leadership Lab"
    )

    # Save
    output_path = '/Users/wendyperdomo/Downloads/GitHub/cwc-executive-lab/Module1-Break-Through-Career-Plateaus.pptx'
    prs.save(output_path)
    print(f"PowerPoint saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
